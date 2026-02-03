from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from zipfile import ZipFile
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont

# === НАСТРОЙКИ НА ПРОЕКТА ===
project_dir = Path("NIKA_TechBlue_Presentation")
project_dir.mkdir(exist_ok=True)
pptx_path = project_dir / "NIKA_TechBlue_Presentation.pptx"
readme_path = project_dir / "README.txt"

# === ЦВЕТОВА ПАЛИТРА ===
dark_blue = RGBColor(0x00, 0x3B, 0x73)
light_blue = RGBColor(0x00, 0xAE, 0xEF)
accent_purple = RGBColor(0xC0, 0x65, 0xFF)

# === СЪЗДАВАНЕ НА ПРЕЗЕНТАЦИЯ ===
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# ФОН и общ footer
def apply_footer(slide):
    txBox = slide.shapes.add_textbox(Inches(0.3), Inches(6.8), Inches(12.5), Inches(0.5))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "👉 ПГЕЕ – гр. Банско | Проект НИКА, 2025"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xE1, 0xF5, 0xFE)
    p.alignment = 2  # Right alignment

# === СЪДЪРЖАНИЕ НА СЛАЙДОВЕ ===
slides_content = [
    ("НИКА – Интелигентна система за видеоанализ в реално време",
     "Проект на ПГЕЕ – гр. Банско\nРъководител: инж. Георги Бориков"),
    ("Какъв проблем решаваме?",
     "Решаваме предизвикателството за автоматично разпознаване и анализ на видеопотоци в реално време."),
    ("Технологична архитектура",
     "Видео вход → AI YOLO модел → Django backend → Vue.js интерфейс."),
    ("AI модул",
     "YOLO и OpenCV за откриване на обекти и оценка на действия."),
    ("Backend",
     "REST API и WebSocket за комуникация в реално време."),
    ("Frontend",
     "Vue.js UI с визуализация на живо видео и статистики."),
    ("Приложения на системата",
     "Индустриални зони, сигурност, контрол на достъп, спортен анализ."),
    ("Екип",
     "Инж. Георги Бориков – Ръководител\nУченици: дизайн, интеграция и логика."),
    ("Заключение",
     "НИКА = Интеграция + Интелект + Иновация."),
    ("Контакти",
     "E-mail: ggborikov@abv.bg  |  ПГЕЕ – гр. Банско"),
]

for title, content in slides_content:
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Заден фон (правоъгълник)
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = light_blue
    line = background.line
    line.color.rgb = light_blue

    # Заглавие
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(12), Inches(1.2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(40)
    p.font.color.rgb = dark_blue

    # Текст
    contentBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.5), Inches(4))
    cf = contentBox.text_frame
    c = cf.add_paragraph()
    c.text = content
    c.font.size = Pt(24)
    c.font.color.rgb = accent_purple

    apply_footer(slide)

# Запис на PowerPoint файла
prs.save(pptx_path)

# === README.txt ===
readme_text = """НИКА – Интелигентна система за видеоанализ в реално време
----------------------------------------------------------
Образователна институция: ПГЕЕ – гр. Банско
Отговорник: инж. Георги Бориков
E-mail: ggborikov@abv.bg
Година: 2025

Footer на всички слайдове:
👉 ПГЕЕ – гр. Банско | Проект НИКА, 2025

Създаден автоматично чрез Python + python-pptx.
"""
with open(readme_path, "w", encoding="utf-8") as f:
    f.write(readme_text)

# === ГЕНЕРИРАНЕ НА ZIP АРХИВ ===
zip_path = Path("NIKA_TechBlue_Presentation.zip")
with ZipFile(zip_path, "w") as zf:
    zf.write(pptx_path, pptx_path.name)
    zf.write(readme_path, readme_path.name)

print(f"✅ Готово: {zip_path.resolve()}")